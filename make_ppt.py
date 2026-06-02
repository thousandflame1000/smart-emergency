# -*- coding: utf-8 -*-
"""產生鄰里守望平台簡報 PPT"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # 完全空白版面

# ── 顏色主題 ──────────────────────────────────
NAVY   = RGBColor(0x1a, 0x52, 0x76)
TEAL   = RGBColor(0x27, 0xAC, 0xB2)
GREEN  = RGBColor(0x27, 0xae, 0x60)
RED    = RGBColor(0xe7, 0x4c, 0x3c)
ORANGE = RGBColor(0xe6, 0x7e, 0x22)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xf0, 0xf4, 0xf8)
DGRAY  = RGBColor(0x44, 0x44, 0x44)


def add_rect(slide, x, y, w, h, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    return shape


def add_text(slide, text, x, y, w, h,
             size=24, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def slide_header(slide, title, subtitle=None, accent=TEAL):
    """頂部色條 + 標題"""
    add_rect(slide, 0, 0, 13.33, 1.4, NAVY)
    add_rect(slide, 0, 1.4, 13.33, 0.08, accent)
    add_text(slide, title, 0.5, 0.15, 12, 0.9,
             size=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.5, 0.85, 12, 0.5,
                 size=16, color=RGBColor(0xaa, 0xd4, 0xe8), align=PP_ALIGN.LEFT)
    add_rect(slide, 0, 7.3, 13.33, 0.2, NAVY)
    add_text(slide, "鄰里守望平台  ·  2026 慈悲科技創新競賽",
             0.3, 7.28, 12, 0.25, size=11,
             color=RGBColor(0x88, 0xaa, 0xcc), align=PP_ALIGN.LEFT)


def bullet_box(slide, lines, x, y, w, h, bg=LGRAY,
               title=None, title_color=NAVY, size=18):
    """帶背景的條列方塊"""
    add_rect(slide, x, y, w, h, bg)
    ty = y + 0.12
    if title:
        add_text(slide, title, x + 0.2, ty, w - 0.3, 0.45,
                 size=20, bold=True, color=title_color)
        ty += 0.45
    for line in lines:
        add_text(slide, line, x + 0.2, ty, w - 0.3, 0.42,
                 size=size, color=DGRAY)
        ty += 0.40


def color_card(slide, x, y, w, h, bg, icon, label, sub, sub_size=16):
    add_rect(slide, x, y, w, h, bg)
    add_text(slide, icon,  x, y + 0.15, w, 0.6,  size=40, align=PP_ALIGN.CENTER)
    add_text(slide, label, x, y + 0.75, w, 0.52, size=22, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, sub,   x, y + 1.25, w, 0.45, size=sub_size,
             color=RGBColor(0xdd, 0xee, 0xff), align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════
# Slide 1 — 封面
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, NAVY)
add_rect(s, 0, 0, 13.33, 0.18, TEAL)
add_rect(s, 0, 7.32, 13.33, 0.18, TEAL)

# 大標題
add_text(s, "🏘️ 鄰里守望平台", 0, 1.6, 13.33, 1.5,
         size=56, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "平時照顧長者，災時守護社區",
         0, 3.1, 13.33, 0.8,
         size=28, color=TEAL, align=PP_ALIGN.CENTER)
add_text(s, "智慧科技  ×  社區互助  ×  LINE 即時通",
         0, 3.9, 13.33, 0.6,
         size=20, color=RGBColor(0x88, 0xbb, 0xdd), align=PP_ALIGN.CENTER)
add_rect(s, 4.5, 4.8, 4.33, 0.06, TEAL)
add_text(s, "2026 第十屆全國慈悲科技創新競賽",
         0, 5.1, 13.33, 0.55,
         size=18, color=RGBColor(0x88, 0xaa, 0xcc), align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════
# Slide 2 — 問題定義
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, LGRAY)
slide_header(s, "問題定義", "台灣獨居長者面臨的困境", accent=RED)

# 4 個統計卡
stats = [
    (RED,    "87 萬人",   "65歲以上獨居長者"),
    (ORANGE, "14 小時",   "緊急狀況到被發現"),
    (NAVY,   "200+ 人",   "每位社工負責長者數"),
    (RGBColor(0x8e,0x44,0xad), "無平台", "災時資源分配混亂"),
]
for i, (bg, num, label) in enumerate(stats):
    cx = 0.4 + i * 3.15
    add_rect(s, cx, 1.7, 2.9, 2.1, bg)
    add_text(s, num,   cx, 1.85, 2.9, 1.0, size=38, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, label, cx, 2.85, 2.9, 0.7, size=16,
             color=RGBColor(0xdd,0xee,0xff), align=PP_ALIGN.CENTER)

# 三大痛點
add_text(s, "三大核心痛點", 0.5, 4.1, 12, 0.5,
         size=22, bold=True, color=NAVY)
pains = [
    ("👴 長者", "身體不適找不到人，求助管道複雜，獨居無人知"),
    ("🙋 志工", "每天電話逐一問候費時，急救知識不足，災時找不到資源"),
    ("🖥️ 管理員", "沒有即時全貌，不知道誰需要幫助、物資在哪裡"),
]
for i, (role, pain) in enumerate(pains):
    cx = 0.4 + i * 4.2
    add_rect(s, cx, 4.6, 4.0, 1.5, WHITE)
    add_text(s, role, cx + 0.15, 4.7, 3.7, 0.45, size=18, bold=True, color=NAVY)
    add_text(s, pain, cx + 0.15, 5.15, 3.7, 0.85, size=15, color=DGRAY)


# ════════════════════════════════════════════════
# Slide 3 — 解法
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, LGRAY)
slide_header(s, "我們的解法", "三合一智慧照護平台", accent=GREEN)

add_text(s, "🟢 日常模式", 0.5, 1.6, 5.8, 0.55, size=22, bold=True, color=GREEN)
add_text(s, "🚨 緊急模式", 7.0, 1.6, 5.8, 0.55, size=22, bold=True, color=RED)

# 日常
daily = [
    "📱 LINE 每日打卡  →  長者按一個按鈕",
    "⏰ 未回應自動升級通知家屬/志工",
    "🤖 AI 急救知識即時回答（RAG）",
    "📊 儀表板即時掌握全局",
]
bullet_box(s, daily, 0.4, 2.15, 5.8, 2.8, WHITE, size=17)

# 緊急
emerg = [
    "🚨 管理員一鍵切換緊急模式",
    "🗺️ 資源地圖即時顯示庇護所/物資站",
    "📦 物資供需自動媒合＋通知志工",
    "📢 廣播 LINE 通知所有用戶",
]
bullet_box(s, emerg, 7.0, 2.15, 5.8, 2.8, WHITE, size=17)

# 中間分隔
add_rect(s, 6.4, 1.6, 0.08, 3.4, RGBColor(0xcc,0xcc,0xcc))
add_text(s, "⟺", 6.1, 2.9, 0.8, 0.8, size=24, color=NAVY, align=PP_ALIGN.CENTER)

# 底部重點
add_rect(s, 0.4, 5.2, 12.4, 0.9, TEAL)
add_text(s, "💡 關鍵創新：用 LINE（長者最熟悉的工具）做照護，零學習成本",
         0.6, 5.28, 12, 0.7, size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════
# Slide 4 — 系統流程
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, LGRAY)
slide_header(s, "系統流程", "三個角色的使用體驗", accent=TEAL)

roles = [
    (GREEN,  "👴 長者", [
        "08:00 收到 LINE 打卡訊息",
        "按 ✅ 我很好  →  打卡完成",
        "不舒服按 🆘 需要幫忙",
        "家屬 + 志工立刻收到通知",
    ]),
    (TEAL,   "🙋 志工", [
        "收到未回應警報（LINE 推播）",
        "LINE 直接問：「老人跌倒怎辦」",
        "AI 3 秒內回答完整 SOP",
        "確認安全後一鍵回報完成",
    ]),
    (NAVY,   "🖥️ 管理員", [
        "儀表板看今日打卡全貌",
        "地圖顯示長者位置（顏色標記）",
        "一鍵切換緊急模式",
        "資源地圖 + 自動媒合物資",
    ]),
]
for i, (bg, role, steps) in enumerate(roles):
    cx = 0.4 + i * 4.2
    add_rect(s, cx, 1.65, 3.9, 0.55, bg)
    add_text(s, role, cx, 1.68, 3.9, 0.5, size=20, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    for j, step in enumerate(steps):
        ty = 2.3 + j * 0.95
        add_rect(s, cx, ty, 3.9, 0.82, WHITE)
        add_rect(s, cx, ty, 0.35, 0.82, bg)
        add_text(s, str(j+1), cx, ty + 0.18, 0.35, 0.45,
                 size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(s, step, cx + 0.42, ty + 0.14, 3.35, 0.55, size=15, color=DGRAY)


# ════════════════════════════════════════════════
# Slide 5 — 技術架構
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, LGRAY)
slide_header(s, "技術架構", "全雲端部署，免費運行", accent=NAVY)

# 架構方塊
add_rect(s, 0.4, 1.6, 12.4, 1.0, NAVY)
add_text(s, "LINE Bot  ←→  FastAPI (Railway 24/7)  ←→  PostgreSQL",
         0.5, 1.75, 12, 0.7, size=20, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)

tech = [
    (TEAL,   "🤖 AI 知識庫 (RAG)",
     ["Gemini embedding-001（3072維向量）",
      "numpy cosine 語義搜尋",
      "知識庫：CPR / 中風 / 地震...",
      "問題 → 向量 → 比對 → 生成回答"]),
    (GREEN,  "⏰ 排程系統",
     ["每天 08:00 發送打卡",
      "每 15 分鐘偵測未回應",
      "每 30 分鐘自動媒合物資",
      "APScheduler 內嵌於後端"]),
    (ORANGE, "📱 LINE Bot",
     ["Flex Message 漂亮按鈕卡片",
      "Rich Menu 固定底部選單",
      "Postback 按鈕回調",
      "自動抓取用戶顯示名稱"]),
]
for i, (bg, title, items) in enumerate(tech):
    cx = 0.4 + i * 4.25
    add_rect(s, cx, 2.85, 3.9, 0.45, bg)
    add_text(s, title, cx, 2.88, 3.9, 0.4,
             size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        add_rect(s, cx, 3.35 + j*0.72, 3.9, 0.65, WHITE)
        add_text(s, "▸ " + item, cx + 0.15, 3.42 + j*0.72, 3.65, 0.55,
                 size=15, color=DGRAY)

# 底部
add_rect(s, 0.4, 6.55, 12.4, 0.65, RGBColor(0xe8,0xf4,0xf8))
add_text(s, "🆓 全免費部署  ·  Railway + PostgreSQL + Gemini API + LINE Bot  =  $0/月",
         0.6, 6.6, 12, 0.5, size=18, bold=True,
         color=NAVY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════
# Slide 6 — 實際畫面
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, LGRAY)
slide_header(s, "系統 Demo", "線上可直接操作", accent=TEAL)

screens = [
    (NAVY,  "📊 主儀表板",
     ["即時打卡狀態（已回應/未回應/等待）",
      "Leaflet 地圖顯示長者位置",
      "物資點標記 + 警報記錄",
      "底部 AI 助手直接問問題"]),
    (GREEN, "⚙️ 管理後台",
     ["長者 / 志工 / 家屬 CRUD",
      "物資管理 + 一鍵自動媒合",
      "知識庫上傳自訂文件",
      "立即發送打卡（Demo 用）"]),
    (TEAL,  "📱 LINE Bot",
     ["底部 Rich Menu 四個快速按鈕",
      "漂亮 Flex Message 打卡卡片",
      "志工直接問問題→AI秒回",
      "自動偵測新用戶並歡迎"]),
]
for i, (bg, title, items) in enumerate(screens):
    cx = 0.4 + i * 4.25
    add_rect(s, cx, 1.65, 3.9, 0.5, bg)
    add_text(s, title, cx, 1.68, 3.9, 0.45, size=19, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    bullet_box(s, ["▸ " + it for it in items],
               cx, 2.2, 3.9, 3.0, WHITE, size=16)

add_rect(s, 0.4, 5.5, 12.4, 0.7, NAVY)
add_text(s, "🌐  https://smart-emergency-production-d744.up.railway.app",
         0.6, 5.57, 12, 0.55, size=20, bold=True,
         color=TEAL, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════
# Slide 7 — 社會影響
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, LGRAY)
slide_header(s, "社會影響", "數字說話", accent=GREEN)

impacts = [
    (GREEN,  "⬆️ 10x",    "照護效率",   "1 位志工可照顧 50+ 長者\n（現況 5-10 位）"),
    (TEAL,   "⬇️ 80%",    "通知延誤",   "自動偵測未回應\n平均 1 小時內通知"),
    (ORANGE, "✅ 即時",    "急救知識",   "SOP 查詢從 Google 10 分鐘\n縮短為 3 秒"),
    (RED,    "🔄 可複製",  "擴散潛力",   "開源部署，任何社區\n里辦公室都能使用"),
]
for i, (bg, num, label, desc) in enumerate(impacts):
    cx = 0.4 + i * 3.15
    add_rect(s, cx, 1.65, 2.9, 3.5, WHITE)
    add_rect(s, cx, 1.65, 2.9, 0.55, bg)
    add_text(s, label, cx, 1.67, 2.9, 0.5, size=17, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, num, cx, 2.3, 2.9, 1.1, size=42, bold=True,
             color=bg, align=PP_ALIGN.CENTER)
    add_text(s, desc, cx + 0.15, 3.5, 2.65, 1.5, size=16, color=DGRAY)

# 願景
add_rect(s, 0.4, 5.5, 12.4, 1.3, NAVY)
add_text(s, "「科技的溫度，在於讓最脆弱的人感受到最即時的關懷。」",
         0.6, 5.6, 12, 0.65, size=22, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "鄰里守望平台  ——  數位時代的守望相助",
         0.6, 6.25, 12, 0.45, size=17,
         color=TEAL, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════
# Slide 8 — 結語 + QR
# ════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, NAVY)
add_rect(s, 0, 0, 13.33, 0.15, TEAL)
add_rect(s, 0, 7.35, 13.33, 0.15, TEAL)

add_text(s, "🏘️ 鄰里守望平台", 0, 0.8, 13.33, 1.0,
         size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "平時照顧長者，災時守護社區",
         0, 1.8, 13.33, 0.7, size=26, color=TEAL, align=PP_ALIGN.CENTER)

add_rect(s, 4.5, 2.65, 4.33, 0.06, TEAL)

links = [
    ("🌐 線上 Demo",  "https://smart-emergency-production-d744.up.railway.app"),
    ("⚙️ 管理後台",   "https://smart-emergency-production-d744.up.railway.app/admin"),
    ("📖 API 文件",   "https://smart-emergency-production-d744.up.railway.app/docs"),
    ("💻 GitHub",    "https://github.com/thousandflame1000/smart-emergency"),
]
for i, (label, url) in enumerate(links):
    cy = 2.9 + i * 0.72
    add_rect(s, 2.5, cy, 8.33, 0.6, RGBColor(0x1f, 0x61, 0x8d))
    add_text(s, label, 2.7, cy + 0.08, 2.5, 0.45, size=16, bold=True, color=TEAL)
    add_text(s, url,   5.0, cy + 0.1,  5.7, 0.42, size=15,
             color=RGBColor(0x88, 0xcc, 0xee))

add_text(s, "感謝聆聽  🙏", 0, 6.0, 13.33, 0.8,
         size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ── 儲存 ──────────────────────────────────────
out = r"D:\Smart_Emergency\presentation.pptx"
prs.save(out)
print("DONE:", out)
